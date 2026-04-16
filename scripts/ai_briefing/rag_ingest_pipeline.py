"""RAG 인덱싱 파이프라인 스캐폴드 (MVP)

- 공통 메타 스키마(rag_documents, rag_chunks, rag_index_jobs) 기준
- 전 도메인 원천 데이터를 문서/청크로 변환
- 임베딩 계산은 후속 단계에서 연결
"""

from __future__ import annotations

from dataclasses import dataclass
from hashlib import sha256
from pathlib import Path
from typing import Iterable, List, Sequence
import json
import re


@dataclass
class DocumentRecord:
    source_type: str
    source_domain: str
    source_id: str
    source_sub_id: str
    title: str
    body_text: str
    summary_text: str
    route_hint: str
    menu_code: str
    page_key: str
    entity_type: str
    owner_dept: str = ""
    security_level: str = "internal"
    permission_scope: str = ""
    tags_json: str = "[]"
    metadata_json: str = "{}"

    @property
    def content_hash(self) -> str:
        payload = "\n".join([self.title, self.body_text, self.summary_text])
        return sha256(payload.encode("utf-8")).hexdigest()


@dataclass
class ChunkRecord:
    chunk_index: int
    chunk_text: str

    @property
    def chunk_hash(self) -> str:
        return sha256(self.chunk_text.encode("utf-8")).hexdigest()

    @property
    def token_count(self) -> int:
        # MVP: 대략적인 토큰 수 추정치
        return max(1, len(self.chunk_text.split()))


def normalize_text(text: str) -> str:
    cleaned = re.sub(r"\s+", " ", (text or "").strip())
    return cleaned


def chunk_text(text: str, chunk_size: int = 420, overlap: int = 60) -> List[ChunkRecord]:
    source = normalize_text(text)
    if not source:
        return []

    words = source.split(" ")
    if not words:
        return []

    chunks: List[ChunkRecord] = []
    step = max(1, chunk_size - overlap)
    idx = 0
    chunk_index = 0

    while idx < len(words):
        piece = words[idx: idx + chunk_size]
        if not piece:
            break
        chunks.append(ChunkRecord(chunk_index=chunk_index, chunk_text=" ".join(piece)))
        chunk_index += 1
        idx += step

    return chunks


def build_document_record(
    *,
    source_domain: str,
    source_name: str,
    source_type: str,
    source_id: str,
    title_parts: Sequence[str],
    body_parts: Sequence[str],
    route_hint: str,
    menu_code: str,
    page_key: str,
    entity_type: str,
    metadata: dict | None = None,
) -> DocumentRecord:
    title = normalize_text(" ".join([p for p in title_parts if p])) or source_id
    body_text = normalize_text("\n".join([p for p in body_parts if p]))
    summary_text = normalize_text(body_text[:300])

    return DocumentRecord(
        source_type=source_type,
        source_domain=source_domain,
        source_id=source_id,
        source_sub_id="",
        title=title,
        body_text=body_text,
        summary_text=summary_text,
        route_hint=route_hint,
        menu_code=menu_code,
        page_key=page_key,
        entity_type=entity_type,
        metadata_json=json.dumps(metadata or {}, ensure_ascii=False),
    )


def iter_attachment_documents() -> Iterable[DocumentRecord]:
    """Deprecated: use iter_attachment_documents_from_paths()."""
    return []


def _extract_text_from_pdf(file_path: Path) -> str:
    try:
        from pypdf import PdfReader
    except Exception as exc:
        raise RuntimeError("pypdf 패키지가 필요합니다.") from exc

    pages: List[str] = []
    reader = PdfReader(str(file_path))
    for page in reader.pages:
        pages.append(page.extract_text() or "")
    return "\n".join(pages)


def _extract_text_from_docx(file_path: Path) -> str:
    try:
        from docx import Document
    except Exception as exc:
        raise RuntimeError("python-docx 패키지가 필요합니다.") from exc

    doc = Document(str(file_path))
    lines = [p.text for p in doc.paragraphs if (p.text or "").strip()]
    return "\n".join(lines)


def _extract_text_from_xlsx(file_path: Path) -> str:
    try:
        from openpyxl import load_workbook
    except Exception as exc:
        raise RuntimeError("openpyxl 패키지가 필요합니다.") from exc

    wb = load_workbook(str(file_path), data_only=True, read_only=True)
    rows: List[str] = []
    for sheet in wb.worksheets:
        rows.append(f"[시트] {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            values = [str(v).strip() for v in row if v is not None and str(v).strip()]
            if values:
                rows.append(" | ".join(values))
    return "\n".join(rows)


def _extract_text_from_pptx(file_path: Path) -> str:
    try:
        from pptx import Presentation
    except Exception as exc:
        raise RuntimeError("python-pptx 패키지가 필요합니다.") from exc

    prs = Presentation(str(file_path))
    lines: List[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        lines.append(f"[슬라이드] {idx}")
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if (text or "").strip():
                lines.append(text.strip())
    return "\n".join(lines)


def extract_text_from_file(file_path: str | Path) -> str:
    path = Path(file_path)
    suffix = path.suffix.lower()

    if suffix == ".pdf":
        return _extract_text_from_pdf(path)
    if suffix == ".docx":
        return _extract_text_from_docx(path)
    if suffix in (".xlsx", ".xlsm"):
        return _extract_text_from_xlsx(path)
    if suffix == ".pptx":
        return _extract_text_from_pptx(path)

    # 텍스트 계열 파일은 UTF-8 우선으로 읽고 실패 시 무시한다.
    if suffix in (".txt", ".md", ".csv", ".log", ".json", ".html"):
        return path.read_text(encoding="utf-8", errors="ignore")

    raise ValueError(f"지원하지 않는 파일 형식입니다: {path.name}")


def iter_attachment_documents_from_paths(
    attachment_paths: Sequence[str],
    *,
    source_domain: str,
    source_id_prefix: str,
    route_hint: str = "",
    menu_code: str = "",
    page_key: str = "",
    entity_type: str = "attachment",
) -> Iterable[DocumentRecord]:
    """첨부파일 경로 목록을 공통 DocumentRecord로 변환한다."""
    for idx, item in enumerate(attachment_paths):
        path = Path(item)
        if not path.exists() or not path.is_file():
            continue

        body_text = normalize_text(extract_text_from_file(path))
        if not body_text:
            continue

        source_id = f"{source_id_prefix}:{idx + 1}"
        yield DocumentRecord(
            source_type="attachment",
            source_domain=source_domain,
            source_id=source_id,
            source_sub_id=path.name,
            title=path.stem,
            body_text=body_text,
            summary_text=normalize_text(body_text[:300]),
            route_hint=route_hint,
            menu_code=menu_code,
            page_key=page_key,
            entity_type=entity_type,
            metadata_json=json.dumps({"file_name": path.name, "file_path": str(path)}, ensure_ascii=False),
        )


def enqueue_reindex_job(source_domain: str, source_id: str, action: str = "upsert") -> dict:
    """큐 시스템 연동 전까지 사용하는 작업 페이로드 스텁."""
    return {
        "source_type": "db_row",
        "source_domain": source_domain,
        "source_id": source_id,
        "source_sub_id": "",
        "action": action,
        "status": "pending",
        "priority": 100,
    }


def build_attachment_job_payload(
    *,
    source_domain: str,
    source_id: str,
    attachment_paths: Sequence[str],
    route_hint: str = "",
    menu_code: str = "",
    page_key: str = "",
    entity_type: str = "attachment",
) -> dict:
    """rag_index_jobs.payload_json에 저장할 첨부 인덱싱 페이로드 생성."""
    return {
        "source_type": "attachment",
        "source_domain": source_domain,
        "source_id": source_id,
        "source_sub_id": "",
        "action": "upsert",
        "status": "pending",
        "priority": 80,
        "payload_json": json.dumps(
            {
                "attachment_paths": list(attachment_paths),
                "route_hint": route_hint,
                "menu_code": menu_code,
                "page_key": page_key,
                "entity_type": entity_type,
            },
            ensure_ascii=False,
        ),
    }


if __name__ == "__main__":
    # 간단 동작 예시
    sample_doc = build_document_record(
        source_domain="인사이트",
        source_name="insight_blog",
        source_type="db_row",
        source_id="blog:1",
        title_parts=["RAG 적용 가이드"],
        body_parts=["RAG 메타데이터 표준화", "권한 필터", "근거 기반 응답"],
        route_hint="/p/insight_blog_it_detail?id=1",
        menu_code="insight",
        page_key="insight_blog_it",
        entity_type="blog",
    )
    sample_chunks = chunk_text(sample_doc.body_text)
    print("doc_hash=", sample_doc.content_hash)
    print("chunks=", len(sample_chunks))
